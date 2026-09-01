"""Build the client-facing Quality Lifecycle Studio demonstration deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Quality_Lifecycle_Studio_Client_Demo.pptx"

BG = RGBColor(6, 8, 15)
SURFACE = RGBColor(16, 23, 40)
SURFACE_2 = RGBColor(22, 31, 51)
TEXT = RGBColor(243, 247, 255)
MUTED = RGBColor(150, 163, 184)
MINT = RGBColor(126, 249, 214)
CYAN = RGBColor(99, 217, 255)
VIOLET = RGBColor(167, 139, 250)
PINK = RGBColor(255, 124, 155)
LINE = RGBColor(53, 65, 88)


def add_text(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=SURFACE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def base_slide(prs, section, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_text(slide, "QUALITY LIFECYCLE STUDIO", 0.55, 0.22, 2.8, 0.3, 9, MINT, True)
    add_text(slide, section.upper(), 9.2, 0.22, 3.4, 0.3, 8, MUTED, False,
             align=PP_ALIGN.RIGHT)
    add_text(slide, f"{number:02d}", 12.2, 7.08, 0.55, 0.22, 8, MUTED,
             align=PP_ALIGN.RIGHT)
    return slide


def title(slide, heading, subtitle=None):
    add_text(slide, heading, 0.58, 0.72, 12.0, 0.7, 27, TEXT, True)
    if subtitle:
        add_text(slide, subtitle, 0.6, 1.4, 11.8, 0.48, 12, MUTED)


def bullet_list(slide, items, x, y, w, h, size=16, accent=MINT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(12)
        paragraph.level = 0
        paragraph.text = f"•  {item}"
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x - 0.18), Inches(y), Inches(0.035), Inches(h)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    return box


def card(slide, heading, body, x, y, w, h, accent=MINT, kicker=None):
    rect(slide, x, y, w, h)
    if kicker:
        add_text(slide, kicker.upper(), x + 0.22, y + 0.15, w - 0.44, 0.2, 8, accent, True)
        heading_y = y + 0.42
    else:
        heading_y = y + 0.22
    add_text(slide, heading, x + 0.22, heading_y, w - 0.44, 0.4, 15, TEXT, True)
    add_text(slide, body, x + 0.22, heading_y + 0.48, w - 0.44, h - 0.85, 10.5, MUTED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Quality Lifecycle Studio — Client Demonstration"
    prs.core_properties.subject = "GitHub Copilot-powered test design POC"
    prs.core_properties.author = "Quality Lifecycle Studio"

    # 1 — Cover
    slide = base_slide(prs, "Client demonstration", 1)
    add_text(slide, "From requirements to\nautomation-ready tests", 0.72, 1.28, 8.6, 1.7,
             34, TEXT, True)
    add_text(slide, "A governed GitHub Copilot test-design proof of concept", 0.76, 3.12,
             7.8, 0.5, 16, MINT, True)
    add_text(slide, "Manual test steps  ·  SpecFlow BDD  ·  Automation/manual feasibility",
             0.76, 3.72, 8.8, 0.35, 12, MUTED)
    rect(slide, 9.55, 1.2, 2.7, 3.7, SURFACE_2, CYAN)
    add_text(slide, "⌁", 10.45, 1.63, 0.9, 0.9, 40, MINT, True, align=PP_ALIGN.CENTER)
    add_text(slide, "GITHUB COPILOT", 9.85, 2.75, 2.1, 0.3, 10, MINT, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Sole AI runtime", 9.85, 3.2, 2.1, 0.3, 13, TEXT, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Local application\nHuman-reviewed output", 9.85, 3.72, 2.1, 0.75,
             10, MUTED, align=PP_ALIGN.CENTER)

    # 2 — Problem and value
    slide = base_slide(prs, "Business value", 2)
    title(slide, "Accelerate the first test-design pass", "A focused assistant for quality engineers—not autonomous QA")
    card(slide, "The challenge", "Requirements must be translated into consistent scenarios before automation starts. This is repetitive, variable, and review intensive.", 0.6, 2.0, 3.8, 3.6, PINK, "Today")
    card(slide, "The intervention", "Generate 4–5 risk-focused scenarios with steps, expected results, SpecFlow Gherkin, synthetic data, and feasibility decisions.", 4.75, 2.0, 3.8, 3.6, CYAN, "POC")
    card(slide, "The outcome", "Faster initial design, clearer automation candidates, explicit manual coverage, and controlled export or Jira publication.", 8.9, 2.0, 3.8, 3.6, MINT, "Value")

    # 3 — Scope
    slide = base_slide(prs, "Solution overview", 3)
    title(slide, "What the POC does", "A concise, review-first workflow")
    bullet_list(slide, [
        "Transforms a story, feature description, or acceptance criteria into structured tests",
        "Produces manual test steps or copy-ready SpecFlow BDD",
        "Uses Scenario Outline and Examples for meaningful data-driven coverage",
        "Segregates automation-feasible and manual scenarios with a rationale",
        "Exports JSON/CSV and publishes selected cases to Jira only on explicit action",
    ], 0.85, 2.0, 7.1, 3.9, 15)
    card(slide, "Designed for review", "Generated output is a draft. The tester owns coverage, domain accuracy, automation implementation, and execution evidence.", 8.45, 2.05, 3.9, 1.75, VIOLET, "Control")
    card(slide, "Predictable POC", "4–5 high-level scenarios keep latency and Copilot premium-request use practical.", 8.45, 4.05, 3.9, 1.55, CYAN, "Scope")

    # 4 — Architecture pipeline
    slide = base_slide(prs, "Architecture", 4)
    title(slide, "Live execution pipeline", "Local orchestration with one approved hosted AI runtime")
    nodes = [
        ("01", "Browser UI", "Story + format"),
        ("02", "FastAPI", "Local Uvicorn"),
        ("03", "Agent service", "Fail-closed registry"),
        ("04", "Copilot SDK", "Authenticated CLI"),
        ("05", "GitHub Copilot", "Organization service"),
        ("06", "Quality gate", "Validate + render"),
    ]
    for i, (num, name, detail) in enumerate(nodes):
        x = 0.45 + i * 2.12
        rect(slide, x, 2.28, 1.72, 1.72, SURFACE, MINT if i in (0, 5) else LINE)
        add_text(slide, num, x + 0.16, 2.45, 0.5, 0.2, 8, MINT, True)
        add_text(slide, name, x + 0.16, 2.83, 1.42, 0.38, 13, TEXT, True)
        add_text(slide, detail, x + 0.16, 3.31, 1.42, 0.38, 9, MUTED)
        if i < len(nodes) - 1:
            add_text(slide, "→", x + 1.77, 2.93, 0.34, 0.35, 18, CYAN, True,
                     align=PP_ALIGN.CENTER)
    add_text(slide, "LOCAL EXECUTION", 0.58, 4.42, 5.75, 0.25, 9, CYAN, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "GITHUB-HOSTED", 8.92, 4.42, 1.72, 0.25, 9, VIOLET, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "LOCAL VALIDATION", 10.98, 4.42, 1.72, 0.25, 9, MINT, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Only the restricted generation request leaves the local application. Jira is a separate, explicit workflow.", 1.0, 5.25, 11.3, 0.5, 12, MUTED, align=PP_ALIGN.CENTER)

    # 5 — LLM / agent / orchestration
    slide = base_slide(prs, "AI system", 5)
    title(slide, "Three distinct layers", "Clear accountability for model, behavior, and control flow")
    card(slide, "LLM", "Selection: organization-default\nGateway: GitHub Copilot\nEntitlement: signed-in organization user\nNo direct provider API key", 0.6, 2.0, 3.8, 3.75, CYAN, "Model")
    card(slide, "Test Designer", "Senior QA role\nNormal + SpecFlow BDD\n4–5 structured scenarios\nAutomation/manual feasibility\nSynthetic data only", 4.75, 2.0, 3.8, 3.75, MINT, "Agent")
    card(slide, "Local pipeline", "FastAPI + Uvicorn\nTestGenerationService\nAgentRegistry policy gate\nPydantic validation\nDeduplication and cap", 8.9, 2.0, 3.8, 3.75, VIOLET, "Orchestration")

    # 6 — Security
    slide = base_slide(prs, "Governance", 6)
    title(slide, "Copilot-only and fail closed", "Controls are implemented in code, not only documented")
    bullet_list(slide, [
        "AgentRegistry rejects every runtime ID except github-copilot",
        "No OpenAI, Anthropic-direct, Ollama, BYOK, or fallback adapter",
        "Tools, MCP servers, skills, memory, file writes, hooks, and Git operations disabled",
        "Model output treated as untrusted and validated with Pydantic",
        "External calls are mocked in CI; automated tests consume no Copilot requests",
        "Jira publication is explicit, selective, and least privileged",
    ], 0.85, 1.92, 8.05, 4.45, 14)
    rect(slide, 9.45, 2.08, 2.85, 2.85, SURFACE_2, MINT)
    add_text(slide, "ONLY", 9.82, 2.55, 2.1, 0.4, 12, MUTED, True, align=PP_ALIGN.CENTER)
    add_text(slide, "github-copilot", 9.65, 3.2, 2.45, 0.45, 20, MINT, True,
             font="Aptos Mono", align=PP_ALIGN.CENTER)
    add_text(slide, "approved runtime", 9.82, 3.92, 2.1, 0.3, 11, TEXT,
             align=PP_ALIGN.CENTER)

    # 7 — Normal demo
    slide = base_slide(prs, "Live demo", 7)
    title(slide, "Demo 1 — Manual test scenarios", "Password-reset story with time-limited links")
    card(slide, "Input", "A registered customer resets a password using a 15-minute email link. Used links cannot be reused; password policy applies; successful reset invalidates sessions.", 0.6, 1.9, 4.15, 3.9, CYAN, "Story")
    card(slide, "Generated structure", "Category + priority\nAutomation/manual classification\nFeasibility rationale\nPreconditions\n1–2 concise steps\nObservable expected results", 5.0, 1.9, 3.25, 3.9, MINT, "Output")
    card(slide, "Presenter focus", "Show the small risk-based suite, point out classification chips, inspect rationale, then export JSON or CSV.", 8.5, 1.9, 4.15, 3.9, VIOLET, "Talk track")

    # 8 — BDD demo
    slide = base_slide(prs, "Live demo", 8)
    title(slide, "Demo 2 — SpecFlow-ready BDD", "Copy one scenario or an entire .feature file")
    rect(slide, 0.62, 1.86, 7.15, 4.65, RGBColor(8, 12, 20), LINE)
    gherkin = (
        "Scenario Outline: Reject an invalid reset attempt\n"
        "  Given a reset link is <link_state>\n"
        "  When the customer submits <new_password>\n"
        "  Then the reset is rejected with <message>\n\n"
        "Examples:\n"
        "  | link_state | new_password | message      |\n"
        "  | expired    | Valid#Pass12 | Link expired |\n"
        "  | used       | Valid#Pass12 | Link used    |\n"
        "  | valid      | short        | Invalid pwd  |"
    )
    add_text(slide, gherkin, 0.9, 2.14, 6.6, 3.95, 13, MINT, font="Aptos Mono")
    card(slide, "Automation hand-off", "Scenario for a single flow\nScenario Outline for repeated data\nExamples tables with placeholders\nReusable, implementation-neutral steps\nCopy Scenario / Copy BDD Feature", 8.15, 1.9, 4.5, 4.55, CYAN, "SpecFlow")

    # 9 — Feasibility
    slide = base_slide(prs, "Test strategy", 9)
    title(slide, "Automation and manual coverage", "The agent explains—not merely labels—the execution decision")
    card(slide, "Automation", "Repeatable and deterministic\nStable UI/API/system interface\nObservable assertions\nHigh rerun value\nExamples: authentication rules, link states, session invalidation", 0.75, 1.95, 5.7, 3.95, MINT, "At least 2")
    card(slide, "Manual", "Exploratory behavior\nSubjective usability or visual clarity\nHuman accessibility experience\nPhysical/CAPTCHA/biometric dependency\nPoor automation ROI", 6.88, 1.95, 5.7, 3.95, VIOLET, "At least 2")
    add_text(slide, "Guardrail: complexity alone is not a valid reason to classify a test as manual.", 1.0, 6.25, 11.3, 0.35, 12, CYAN, True, align=PP_ALIGN.CENTER)

    # 10 — Integrations
    slide = base_slide(prs, "Workflow", 10)
    title(slide, "Controlled downstream integration", "Human approval remains in the loop")
    card(slide, "Copy", "Copy an individual Gherkin scenario or a complete SpecFlow Feature block to the clipboard.", 0.6, 2.0, 2.8, 3.4, CYAN, "01")
    card(slide, "Export", "Download the validated suite as CSV or JSON for review, evidence, or downstream transformation.", 3.7, 2.0, 2.8, 3.4, MINT, "02")
    card(slide, "Select", "Choose which scenarios are eligible for publication rather than sending the complete draft.", 6.8, 2.0, 2.8, 3.4, VIOLET, "03")
    card(slide, "Jira", "Attach selected cases and add a summary only after the tester explicitly initiates publishing.", 9.9, 2.0, 2.8, 3.4, PINK, "04")

    # 11 — Limitations
    slide = base_slide(prs, "POC boundaries", 11)
    title(slide, "What this demonstration does not claim", "Transparent limitations support responsible adoption")
    bullet_list(slide, [
        "Generated scenarios are suggestions, not proof of complete coverage",
        "The POC does not execute tests or generate SpecFlow step-definition code",
        "Gherkin requires alignment with the client's domain vocabulary and shared bindings",
        "No application user management, RBAC, audit database, or persistent suite history",
        "Copilot model availability, latency, limits, and billing follow organization policy",
        "Current Jira integration attaches CSV; it is not an Xray or Zephyr adapter",
    ], 0.9, 1.85, 11.2, 4.75, 14, PINK)

    # 12 — Demo runbook
    slide = base_slide(prs, "Presenter runbook", 12)
    title(slide, "10-minute client walkthrough", "Keep the narrative focused on value, controls, and measurable next steps")
    steps = [
        ("00–01", "Problem and value"), ("01–02", "Architecture and governance"),
        ("02–04", "Normal scenario generation"), ("04–07", "SpecFlow BDD + copy"),
        ("07–08", "Automation/manual mix"), ("08–09", "Export and Jira"),
        ("09–10", "Limitations and pilot proposal"),
    ]
    for i, (time, label) in enumerate(steps):
        y = 1.78 + i * 0.67
        add_text(slide, time, 0.9, y, 1.1, 0.32, 11, MINT, True, font="Aptos Mono")
        add_text(slide, label, 2.15, y, 5.3, 0.32, 14, TEXT, True)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.7), Inches(y + 0.12), Inches(4.2), Inches(0.02))
        line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.fill.background()

    # 13 — Next steps
    slide = base_slide(prs, "Recommendation", 13)
    title(slide, "Move from POC to a measured pilot", "Validate usefulness, edit effort, latency, governance, and adoption")
    cards = [
        ("1", "Select workflows", "Choose 2–3 representative, non-sensitive client journeys."),
        ("2", "Define measures", "Track design time, usefulness, edit rate, and review defects."),
        ("3", "Align standards", "Map Gherkin to the client's SpecFlow vocabulary and bindings."),
        ("4", "Review controls", "Agree Copilot policy, model, budget, privacy, logs, and deployment."),
        ("5", "Run pilot", "Time-box tester usage and capture structured feedback."),
    ]
    for i, (num, heading, body) in enumerate(cards):
        x = 0.47 + i * 2.55
        card(slide, heading, body, x, 2.05, 2.3, 3.65, MINT if i < 3 else CYAN, num)

    # 14 — Close
    slide = base_slide(prs, "Discussion", 14)
    add_text(slide, "Accelerate design.\nPreserve judgment.", 0.75, 1.35, 8.4, 1.5, 38, TEXT, True)
    add_text(slide, "A governed path from requirements to an initial automation-ready test design using the organization’s GitHub Copilot capability.", 0.8, 3.22, 8.5, 0.8, 16, MUTED)
    rect(slide, 9.65, 1.4, 2.65, 3.25, SURFACE_2, MINT)
    add_text(slide, "QUESTIONS", 9.95, 2.2, 2.05, 0.35, 11, MINT, True, align=PP_ALIGN.CENTER)
    add_text(slide, "&", 10.58, 2.8, 0.8, 0.55, 28, TEXT, True, align=PP_ALIGN.CENTER)
    add_text(slide, "NEXT STEPS", 9.95, 3.58, 2.05, 0.35, 11, CYAN, True, align=PP_ALIGN.CENTER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
